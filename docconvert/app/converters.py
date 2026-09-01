"""Format specific converters for the LLMInt document converter service.

Every converter turns a binary payload into a :class:`Document`, i.e. a list of
logical blocks with a human readable source reference. Keeping the structure
(headings, sheets, slides, tables) instead of a flat text dump makes the
resulting RAG chunks far more useful.
"""

from __future__ import annotations

import csv
import io
import json
import re
import zipfile
from pathlib import Path

from .chunking import Document, normalize_text


class UnsupportedFormatError(Exception):
    """Raised when no converter can handle the given file."""


class ConversionError(Exception):
    """Raised when a converter fails on an otherwise supported file."""


# ── Format registry ───────────────────────────────────────────────────────────
# extension -> (label, mime types accepted for this extension)
SUPPORTED_FORMATS = {
    ".docx": ("Word-Dokument", ["application/vnd.openxmlformats-officedocument.wordprocessingml.document"]),
    ".xlsx": ("Excel-Arbeitsmappe", ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"]),
    ".xlsm": ("Excel-Arbeitsmappe (Makros)", ["application/vnd.ms-excel.sheet.macroEnabled.12"]),
    ".xls": ("Excel-Arbeitsmappe (97-2003)", ["application/vnd.ms-excel"]),
    ".pptx": ("PowerPoint-Präsentation", ["application/vnd.openxmlformats-officedocument.presentationml.presentation"]),
    ".odt": ("OpenDocument-Text", ["application/vnd.oasis.opendocument.text"]),
    ".ods": ("OpenDocument-Tabelle", ["application/vnd.oasis.opendocument.spreadsheet"]),
    ".odp": ("OpenDocument-Präsentation", ["application/vnd.oasis.opendocument.presentation"]),
    ".rtf": ("Rich Text", ["application/rtf", "text/rtf"]),
    ".csv": ("CSV-Tabelle", ["text/csv"]),
    ".tsv": ("TSV-Tabelle", ["text/tab-separated-values"]),
    ".txt": ("Textdatei", ["text/plain"]),
    ".md": ("Markdown", ["text/markdown", "text/plain"]),
    ".markdown": ("Markdown", ["text/markdown", "text/plain"]),
    ".json": ("JSON-Datei", ["application/json"]),
    ".xml": ("XML-Datei", ["application/xml", "text/xml"]),
    ".yaml": ("YAML-Datei", ["application/yaml", "text/yaml"]),
    ".yml": ("YAML-Datei", ["application/yaml", "text/yaml"]),
    ".html": ("HTML-Dokument", ["text/html"]),
    ".htm": ("HTML-Dokument", ["text/html"]),
    ".log": ("Logdatei", ["text/plain"]),
    ".ini": ("Konfigurationsdatei", ["text/plain"]),
    ".conf": ("Konfigurationsdatei", ["text/plain"]),
}

# Rows of a spreadsheet that go into one chunkable block.
SHEET_ROWS_PER_BLOCK = 40
# Hard limit for cells rendered per sheet, protects against runaway workbooks.
MAX_SHEET_CELLS = 200_000


def detect_extension(filename: str, mime_type: str = "") -> str:
    """Resolve the converter key for a file name, falling back to the MIME type."""

    ext = Path(filename or "").suffix.lower()
    if ext in SUPPORTED_FORMATS:
        return ext

    mime = (mime_type or "").split(";")[0].strip().lower()
    for candidate, (_label, mimes) in SUPPORTED_FORMATS.items():
        if mime and mime in mimes:
            return candidate

    if mime.startswith("text/"):
        return ".txt"

    raise UnsupportedFormatError(
        f"Format wird nicht unterstützt: {ext or mime or 'unbekannt'}"
    )


# ── Helpers ───────────────────────────────────────────────────────────────────


def _decode(content: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def _cell_to_str(value) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()


def _rows_to_markdown(rows: list) -> str:
    """Render a list of row lists as a Markdown table (first row = header)."""

    if not rows:
        return ""

    width = max(len(row) for row in rows)
    padded = [[_cell_to_str(cell) for cell in row] + [""] * (width - len(row)) for row in rows]

    header = padded[0]
    if not any(header):
        header = [f"Spalte {i + 1}" for i in range(width)]
        body = padded
    else:
        body = padded[1:]

    lines = ["| " + " | ".join(header) + " |", "|" + "|".join(["---"] * width) + "|"]
    for row in body:
        if not any(cell for cell in row):
            continue
        lines.append("| " + " | ".join(cell.replace("|", "\\|") for cell in row) + " |")
    return "\n".join(lines)


def _table_blocks(document: Document, rows: list, source_ref: str, group: str,
                  label: str = "") -> None:
    """Add a (possibly large) table to the document in row windows."""

    if not rows:
        return

    header = rows[0]
    body = rows[1:]

    if not body:
        document.add(_rows_to_markdown(rows), source_ref, group)
        return

    for start in range(0, len(body), SHEET_ROWS_PER_BLOCK):
        window = body[start:start + SHEET_ROWS_PER_BLOCK]
        ref = source_ref
        if len(body) > SHEET_ROWS_PER_BLOCK:
            ref = f"{source_ref}, Zeilen {start + 1}–{start + len(window)}"
        text = _rows_to_markdown([header] + window)
        if label:
            text = f"{label}\n{text}"
        document.add(text, ref, f"{group}#{start // SHEET_ROWS_PER_BLOCK}")


# ── Word ──────────────────────────────────────────────────────────────────────


def convert_docx(content: bytes, filename: str) -> Document:
    import docx  # python-docx
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    try:
        source = docx.Document(io.BytesIO(content))
    except Exception as exc:  # noqa: BLE001 - surfaced to the caller as ConversionError
        raise ConversionError(f"Word-Dokument konnte nicht gelesen werden: {exc}") from exc

    document = Document()
    heading_stack: list = []
    table_index = 0

    def current_ref() -> str:
        return " › ".join(heading_stack) if heading_stack else "Dokument"

    body = source.element.body
    for child in body.iterchildren():
        tag = child.tag.split("}")[-1]

        if tag == "p":
            paragraph = Paragraph(child, source)
            text = paragraph.text.strip()
            if not text:
                continue

            style = (paragraph.style.name if paragraph.style is not None else "") or ""
            match = re.match(r"(?:Heading|Überschrift)\s*(\d+)", style, re.IGNORECASE)
            if match:
                level = max(1, min(9, int(match.group(1))))
                del heading_stack[level - 1:]
                heading_stack.append(text)
                document.add("#" * level + " " + text, current_ref(), current_ref())
            else:
                document.add(text, current_ref(), current_ref())

        elif tag == "tbl":
            table = Table(child, source)
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            table_index += 1
            ref = f"{current_ref()}, Tabelle {table_index}"
            _table_blocks(document, rows, ref, f"table-{table_index}")

    # Headers/footers are skipped on purpose: they repeat on every page and add
    # noise to the embeddings.
    core = source.core_properties
    document.meta = {
        "paragraphs": len(source.paragraphs),
        "tables": table_index,
        "title": (core.title or "").strip(),
        "author": (core.author or "").strip(),
    }
    return document


# ── Spreadsheets ──────────────────────────────────────────────────────────────


def _sheet_rows_xlsx(content: bytes) -> dict:
    import openpyxl

    try:
        workbook = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    except Exception as exc:  # noqa: BLE001
        raise ConversionError(f"Arbeitsmappe konnte nicht gelesen werden: {exc}") from exc

    sheets = {}
    cells = 0
    try:
        for worksheet in workbook.worksheets:
            rows = []
            for row in worksheet.iter_rows(values_only=True):
                values = [_cell_to_str(value) for value in row]
                while values and not values[-1]:
                    values.pop()
                cells += len(values)
                if cells > MAX_SHEET_CELLS:
                    rows.append(["… (Arbeitsmappe gekürzt)"])
                    break
                if any(values):
                    rows.append(values)
            sheets[worksheet.title] = rows
    finally:
        workbook.close()
    return sheets


def _sheet_rows_xls(content: bytes) -> dict:
    import xlrd

    try:
        book = xlrd.open_workbook(file_contents=content)
    except Exception as exc:  # noqa: BLE001
        raise ConversionError(f"Arbeitsmappe (.xls) konnte nicht gelesen werden: {exc}") from exc

    sheets = {}
    for sheet in book.sheets():
        rows = []
        for row_index in range(sheet.nrows):
            values = [_cell_to_str(cell) for cell in sheet.row_values(row_index)]
            while values and not values[-1]:
                values.pop()
            if any(values):
                rows.append(values)
        sheets[sheet.name] = rows
    return sheets


def _sheet_rows_ods(content: bytes) -> dict:
    import pandas as pd

    try:
        frames = pd.read_excel(io.BytesIO(content), sheet_name=None, engine="odf", header=None,
                               dtype=str)
    except Exception as exc:  # noqa: BLE001
        raise ConversionError(f"OpenDocument-Tabelle konnte nicht gelesen werden: {exc}") from exc

    sheets = {}
    for name, frame in frames.items():
        frame = frame.fillna("")
        rows = [[_cell_to_str(value) for value in row] for row in frame.values.tolist()]
        sheets[name] = [row for row in rows if any(row)]
    return sheets


def _spreadsheet_document(sheets: dict, kind: str) -> Document:
    document = Document()
    non_empty = 0

    for name, rows in sheets.items():
        if not rows:
            continue
        non_empty += 1
        ref = f'Blatt "{name}"'
        _table_blocks(document, rows, ref, f"sheet-{name}",
                      label=f"Tabellenblatt: {name} ({len(rows) - 1} Datenzeilen)")

    document.meta = {
        "kind": kind,
        "sheets": list(sheets.keys()),
        "sheet_count": len(sheets),
        "non_empty_sheets": non_empty,
    }
    return document


def convert_xlsx(content: bytes, filename: str) -> Document:
    return _spreadsheet_document(_sheet_rows_xlsx(content), "xlsx")


def convert_xls(content: bytes, filename: str) -> Document:
    return _spreadsheet_document(_sheet_rows_xls(content), "xls")


def convert_ods(content: bytes, filename: str) -> Document:
    return _spreadsheet_document(_sheet_rows_ods(content), "ods")


def convert_csv(content: bytes, filename: str, delimiter: str = "") -> Document:
    text = _decode(content)
    if not delimiter:
        sample = text[:8192]
        try:
            delimiter = csv.Sniffer().sniff(sample, delimiters=",;\t|").delimiter
        except csv.Error:
            delimiter = "\t" if Path(filename or "").suffix.lower() == ".tsv" else ","

    rows = [row for row in csv.reader(io.StringIO(text), delimiter=delimiter) if any(
        cell.strip() for cell in row)]

    document = Document()
    name = Path(filename or "Tabelle").name
    _table_blocks(document, rows, f'Tabelle "{name}"', "csv",
                  label=f"Tabelle: {name} ({max(0, len(rows) - 1)} Datenzeilen)")
    document.meta = {
        "kind": "csv",
        "delimiter": delimiter,
        "rows": max(0, len(rows) - 1),
        "columns": len(rows[0]) if rows else 0,
    }
    return document


# ── Presentations ─────────────────────────────────────────────────────────────


def convert_pptx(content: bytes, filename: str) -> Document:
    from pptx import Presentation

    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as exc:  # noqa: BLE001
        raise ConversionError(f"Präsentation konnte nicht gelesen werden: {exc}") from exc

    document = Document()
    slide_count = 0

    for index, slide in enumerate(presentation.slides, start=1):
        slide_count += 1
        ref = f"Folie {index}"
        parts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                table = _rows_to_markdown(rows)
                if table:
                    parts.append(table)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            parts.append(f"Notizen: {notes}")

        document.add("\n\n".join(parts), ref, f"slide-{index}")

    document.meta = {"kind": "pptx", "slides": slide_count}
    return document


def convert_odp(content: bytes, filename: str) -> Document:
    from odf.opendocument import load
    from odf import text as odf_text
    from odf.draw import Page

    try:
        odf_document = load(io.BytesIO(content))
    except Exception as exc:  # noqa: BLE001
        raise ConversionError(f"OpenDocument-Präsentation konnte nicht gelesen werden: {exc}") from exc

    document = Document()
    pages = odf_document.getElementsByType(Page)

    for index, page in enumerate(pages, start=1):
        parts = []
        for paragraph in page.getElementsByType(odf_text.P):
            value = str(paragraph).strip()
            if value:
                parts.append(value)
        document.add("\n".join(parts), f"Folie {index}", f"slide-{index}")

    document.meta = {"kind": "odp", "slides": len(pages)}
    return document


# ── OpenDocument text ─────────────────────────────────────────────────────────


def convert_odt(content: bytes, filename: str) -> Document:
    from odf.opendocument import load
    from odf import text as odf_text
    from odf.table import Table as OdfTable, TableRow, TableCell

    try:
        odf_document = load(io.BytesIO(content))
    except Exception as exc:  # noqa: BLE001
        raise ConversionError(f"OpenDocument-Text konnte nicht gelesen werden: {exc}") from exc

    document = Document()
    heading_stack: list = []

    def current_ref() -> str:
        return " › ".join(heading_stack) if heading_stack else "Dokument"

    body = odf_document.text
    table_index = 0

    for node in body.childNodes:
        name = getattr(node, "qname", (None, ""))[1]

        if name == "h":
            text = str(node).strip()
            if not text:
                continue
            try:
                level = int(node.getAttribute("outlinelevel") or 1)
            except (TypeError, ValueError):
                level = 1
            level = max(1, min(9, level))
            del heading_stack[level - 1:]
            heading_stack.append(text)
            document.add("#" * level + " " + text, current_ref(), current_ref())

        elif name == "p":
            text = str(node).strip()
            if text:
                document.add(text, current_ref(), current_ref())

        elif name == "table":
            table_index += 1
            rows = []
            for row in node.getElementsByType(TableRow):
                cells = []
                for cell in row.getElementsByType(TableCell):
                    cells.append(" ".join(
                        str(p).strip() for p in cell.getElementsByType(odf_text.P)
                    ).strip())
                if any(cells):
                    rows.append(cells)
            _table_blocks(document, rows, f"{current_ref()}, Tabelle {table_index}",
                          f"table-{table_index}")

    if not document.blocks:
        # Flat fallback for documents with an unusual body structure.
        for paragraph in odf_document.getElementsByType(odf_text.P):
            document.add(str(paragraph).strip(), "Dokument", "Dokument")

    document.meta = {"kind": "odt", "tables": table_index}
    return document


# ── Plain text family ─────────────────────────────────────────────────────────


def convert_markdown(content: bytes, filename: str) -> Document:
    text = normalize_text(_decode(content))
    document = Document()
    heading_stack: list = []
    buffer: list = []

    def current_ref() -> str:
        return " › ".join(heading_stack) if heading_stack else Path(filename or "Dokument").name

    def flush() -> None:
        if buffer:
            document.add("\n".join(buffer), current_ref(), current_ref())
            buffer.clear()

    for line in text.split("\n"):
        heading = re.match(r"^(#{1,6})\s+(.*)$", line)
        if heading:
            flush()
            level = len(heading.group(1))
            title = heading.group(2).strip()
            del heading_stack[level - 1:]
            heading_stack.append(title)
            buffer.append(line)
            continue
        buffer.append(line)

    flush()
    document.meta = {"kind": "markdown", "characters": len(text)}
    return document


def convert_plain(content: bytes, filename: str) -> Document:
    text = normalize_text(_decode(content))
    document = Document()
    ref = Path(filename or "Dokument").name

    for paragraph in re.split(r"\n{2,}", text):
        document.add(paragraph, ref, ref)

    if not document.blocks:
        document.add(text, ref, ref)

    document.meta = {"kind": "text", "characters": len(text)}
    return document


def convert_json(content: bytes, filename: str) -> Document:
    raw = _decode(content)
    try:
        parsed = json.loads(raw)
    except ValueError as exc:
        raise ConversionError(f"JSON konnte nicht gelesen werden: {exc}") from exc

    document = Document()
    ref = Path(filename or "Dokument").name

    if isinstance(parsed, list):
        for index, item in enumerate(parsed, start=1):
            document.add(json.dumps(item, ensure_ascii=False, indent=2),
                         f"{ref}, Eintrag {index}", f"item-{index}")
    elif isinstance(parsed, dict):
        for key, value in parsed.items():
            document.add(f"{key}: " + json.dumps(value, ensure_ascii=False, indent=2),
                         f"{ref} › {key}", f"key-{key}")
    else:
        document.add(json.dumps(parsed, ensure_ascii=False), ref, ref)

    document.meta = {"kind": "json", "root": type(parsed).__name__}
    return document


def convert_html(content: bytes, filename: str) -> Document:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_decode(content), "lxml")
    for tag in soup(["script", "style", "noscript", "template"]):
        tag.decompose()

    document = Document()
    title = (soup.title.get_text(strip=True) if soup.title else "") or Path(filename or "Dokument").name
    heading_stack: list = []

    def current_ref() -> str:
        return " › ".join([title] + heading_stack) if heading_stack else title

    for element in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "pre", "table"]):
        text = element.get_text(" ", strip=True)
        if not text:
            continue

        if element.name.startswith("h") and len(element.name) == 2 and element.name[1].isdigit():
            level = int(element.name[1])
            del heading_stack[level - 1:]
            heading_stack.append(text)
            document.add("#" * level + " " + text, current_ref(), current_ref())
        elif element.name == "table":
            rows = [[cell.get_text(" ", strip=True) for cell in row.find_all(["th", "td"])]
                    for row in element.find_all("tr")]
            rows = [row for row in rows if any(row)]
            _table_blocks(document, rows, current_ref(), f"table-{id(element)}")
        else:
            document.add(text, current_ref(), current_ref())

    document.meta = {"kind": "html", "title": title}
    return document


def convert_xml(content: bytes, filename: str) -> Document:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_decode(content), "xml")
    document = Document()
    ref = Path(filename or "Dokument").name
    root = soup.find()

    children = list(root.find_all(recursive=False)) if root is not None else []
    if children:
        for index, child in enumerate(children, start=1):
            text = child.get_text(" ", strip=True)
            if text:
                document.add(f"<{child.name}> {text}", f"{ref}, {child.name} {index}",
                             f"node-{index}")
    else:
        document.add(soup.get_text(" ", strip=True), ref, ref)

    document.meta = {"kind": "xml", "root": root.name if root is not None else ""}
    return document


def convert_rtf(content: bytes, filename: str) -> Document:
    from striprtf.striprtf import rtf_to_text

    try:
        text = rtf_to_text(_decode(content), errors="ignore")
    except Exception as exc:  # noqa: BLE001
        raise ConversionError(f"RTF konnte nicht gelesen werden: {exc}") from exc

    document = convert_plain(text.encode("utf-8"), filename)
    document.meta = {"kind": "rtf", "characters": len(text)}
    return document


# ── Dispatch ──────────────────────────────────────────────────────────────────

_CONVERTERS = {
    ".docx": convert_docx,
    ".xlsx": convert_xlsx,
    ".xlsm": convert_xlsx,
    ".xls": convert_xls,
    ".pptx": convert_pptx,
    ".odt": convert_odt,
    ".ods": convert_ods,
    ".odp": convert_odp,
    ".rtf": convert_rtf,
    ".csv": convert_csv,
    ".tsv": convert_csv,
    ".json": convert_json,
    ".html": convert_html,
    ".htm": convert_html,
    ".xml": convert_xml,
    ".md": convert_markdown,
    ".markdown": convert_markdown,
    ".txt": convert_plain,
    ".yaml": convert_plain,
    ".yml": convert_plain,
    ".log": convert_plain,
    ".ini": convert_plain,
    ".conf": convert_plain,
}


def convert(content: bytes, filename: str, mime_type: str = "") -> tuple:
    """Convert ``content`` into a :class:`Document`.

    Returns a ``(document, extension)`` tuple. Raises
    :class:`UnsupportedFormatError` or :class:`ConversionError`.
    """

    if not content:
        raise ConversionError("Die Datei ist leer.")

    ext = detect_extension(filename, mime_type)

    # Office Open XML files are ZIP containers; a mismatching extension is a
    # common cause of confusing parser errors, so verify it up front.
    if ext in {".docx", ".xlsx", ".xlsm", ".pptx"} and not zipfile.is_zipfile(io.BytesIO(content)):
        raise ConversionError(
            "Die Datei ist kein gültiges Office-Open-XML-Dokument (ZIP-Container erwartet)."
        )

    converter = _CONVERTERS.get(ext)
    if converter is None:
        raise UnsupportedFormatError(f"Kein Konverter für {ext}.")

    document = converter(content, filename)
    if not document.blocks:
        raise ConversionError("Die Datei enthält keinen auslesbaren Text.")

    document.meta.setdefault("format", ext.lstrip("."))
    document.meta["label"] = SUPPORTED_FORMATS.get(ext, ("Dokument", []))[0]
    return document, ext
